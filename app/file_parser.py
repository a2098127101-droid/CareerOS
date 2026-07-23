from __future__ import annotations

import csv
import io
import json
from pathlib import Path

from docx import Document
from pypdf import PdfReader


def parse_uploaded_file(filename: str, content: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in {".txt", ".md", ".markdown"}:
        return content.decode("utf-8", errors="replace")
    if suffix == ".docx":
        doc = Document(io.BytesIO(content))
        parts: list[str] = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                vals = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                if any(vals):
                    parts.append(" | ".join(vals))
        return "\n".join(parts)
    if suffix == ".pdf":
        reader = PdfReader(io.BytesIO(content))
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        parts: list[str] = []
        for index, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    texts.append(text.strip())
            if texts:
                parts.append(f"[Slide {index}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    if suffix == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"[Sheet: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                vals = ["" if v is None else str(v).strip() for v in row]
                if any(vals):
                    parts.append(" | ".join(vals))
        return "\n".join(parts)
    if suffix == ".csv":
        text = content.decode("utf-8-sig", errors="replace")
        rows = csv.reader(io.StringIO(text))
        return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)
    if suffix == ".json":
        obj = json.loads(content.decode("utf-8-sig", errors="replace"))
        return json.dumps(obj, ensure_ascii=False, indent=2)
    raise ValueError("仅支持 DOCX / PPTX / XLSX / PDF / TXT / MD / CSV / JSON 文件")
